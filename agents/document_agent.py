"""
Document Agent - Proactive Document Analysis
Extracts technologies, policies, and Zero Trust controls from uploaded documents
"""

from typing import Dict, Any, List
import os
from pathlib import Path

from .base_agent import BaseAgent

# Document parsing imports
try:
    from docx import Document as DocxDocument
    import PyPDF2
    from pptx import Presentation
    import pandas as pd
except ImportError:
    pass  # Will be installed in requirements


class DocumentAgent(BaseAgent):
    """
    Analyzes uploaded documents to extract Zero Trust relevant information.
    
    This is NOT just RAG - it's proactive analysis that:
    1. Parses multiple document formats
    2. Extracts structured data (technologies, controls, policies)
    3. Prepares summaries for next agents
    4. Identifies gaps and areas needing assessment
    """
    
    def __init__(self, **kwargs):
        super().__init__(name="DocumentAgent", **kwargs)
        self.supported_formats = ['.pdf', '.docx', '.pptx', '.xlsx', '.csv', '.txt']
    
    async def process(self, state: Dict[str, Any]) -> Dict[str, Any]:
        """
        Main processing: Parse documents and extract ZT-relevant data.
        
        Expected state keys:
            - uploaded_files: List of file paths
            
        Returns updated state with:
            - documents_analyzed: List of analyzed docs
            - extracted_technologies: List of tech stack
            - extracted_controls: List of ZT controls found
            - extracted_policies: List of policies mentioned
            - document_summaries: Dict of file -> summary
        """
        self.logger.info("Starting document analysis")
        
        # Validate input
        self.validate_state(state, ['uploaded_files'])
        
        uploaded_files = state['uploaded_files']
        
        if not uploaded_files:
            self.logger.warning("No files to analyze")
            return self.update_state(state, {
                'documents_analyzed': [],
                'extraction_status': 'no_files'
            })
        
        # Process each document
        all_text = []
        document_summaries = {}
        
        for file_path in uploaded_files:
            self.logger.info(f"Analyzing: {file_path}")
            
            try:
                text = await self._parse_document(file_path)
                all_text.append(text)
                
                # Create summary for this doc
                summary = await self._create_document_summary(file_path, text)
                document_summaries[os.path.basename(file_path)] = summary
                
            except Exception as e:
                self.logger.error(f"Error parsing {file_path}: {str(e)}")
                document_summaries[os.path.basename(file_path)] = {
                    'error': str(e),
                    'status': 'failed'
                }
        
        # Extract structured data from all documents
        combined_text = "\n\n".join(all_text)
        
        extracted_data = await self._extract_zt_elements(combined_text)
        
        # Update state
        updates = {
            'documents_analyzed': len(uploaded_files),
            'document_summaries': document_summaries,
            'extracted_technologies': extracted_data['technologies'],
            'extracted_controls': extracted_data['controls'],
            'extracted_policies': extracted_data['policies'],
            'extraction_status': 'complete',
            'combined_document_text': combined_text[:10000]  # First 10k chars for context
        }
        
        self.logger.info(f"Extracted: {len(extracted_data['technologies'])} technologies, "
                        f"{len(extracted_data['controls'])} controls, "
                        f"{len(extracted_data['policies'])} policies")
        
        return self.update_state(state, updates)
    
    async def _parse_document(self, file_path: str) -> str:
        """Parse document based on file type."""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            return self._parse_pdf(file_path)
        elif ext == '.docx':
            return self._parse_docx(file_path)
        elif ext == '.pptx':
            return self._parse_pptx(file_path)
        elif ext in ['.xlsx', '.csv']:
            return self._parse_spreadsheet(file_path)
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    
    def _parse_pdf(self, file_path: str) -> str:
        """Extract text from PDF."""
        text = []
        
        with open(file_path, 'rb') as f:
            pdf = PyPDF2.PdfReader(f)
            
            for page in pdf.pages:
                text.append(page.extract_text())
        
        return "\n\n".join(text)
    
    def _parse_docx(self, file_path: str) -> str:
        """Extract text from Word document."""
        doc = DocxDocument(file_path)
        
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text.append(cell.text)
        
        return "\n\n".join(text)
    
    def _parse_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint."""
        prs = Presentation(file_path)
        
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text.append(shape.text)
        
        return "\n\n".join(text)
    
    def _parse_spreadsheet(self, file_path: str) -> str:
        """Extract text from Excel/CSV."""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.csv':
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
        
        # Convert to text representation
        return df.to_string()
    
    async def _create_document_summary(
        self,
        file_path: str,
        text: str
    ) -> Dict[str, Any]:
        """Create a summary of the document using Claude."""
        
        # Only send a summary to Claude, not full document (security)
        preview = text[:2000]  # First 2000 chars
        
        system_prompt = """You are analyzing a document for Zero Trust security assessment.
        Create a brief summary focusing on:
        - Document type and purpose
        - Key security-relevant information
        - Technologies or systems mentioned
        - Any Zero Trust controls discussed"""
        
        user_prompt = f"""Document: {os.path.basename(file_path)}

Preview:
{preview}

Provide a concise 3-4 sentence summary."""
        
        try:
            summary_text = await self.call_claude(system_prompt, user_prompt)
            
            return {
                'filename': os.path.basename(file_path),
                'summary': summary_text,
                'length': len(text),
                'status': 'success'
            }
        except Exception as e:
            return {
                'filename': os.path.basename(file_path),
                'error': str(e),
                'status': 'failed'
            }
    
    async def _extract_zt_elements(self, text: str) -> Dict[str, List[str]]:
        """
        Extract Zero Trust elements using Claude.
        
        This is proactive extraction, not RAG retrieval.
        """
        
        # Send structured extraction request
        system_prompt = """You are a Zero Trust security expert extracting information from documents.

Extract and list:
1. Technologies: Cloud platforms, identity systems, networking tools, security products
2. Controls: Specific ZT controls mentioned (MFA, least privilege, encryption, etc.)
3. Policies: Security policies, access policies, data policies referenced

Format as:
TECHNOLOGIES:
- item1
- item2

CONTROLS:
- item1
- item2

POLICIES:
- item1
- item2"""
        
        # Only send bullet points, not full text (security)
        text_sample = text[:5000]  # First 5000 chars
        
        user_prompt = f"""Extract Zero Trust elements from this text:

{text_sample}

Provide structured extraction as specified."""
        
        try:
            response = await self.call_claude(system_prompt, user_prompt)
            
            # Parse response
            technologies = self._extract_list(response, 'TECHNOLOGIES:')
            controls = self._extract_list(response, 'CONTROLS:')
            policies = self._extract_list(response, 'POLICIES:')
            
            return {
                'technologies': technologies,
                'controls': controls,
                'policies': policies
            }
            
        except Exception as e:
            self.logger.error(f"Extraction failed: {str(e)}")
            return {
                'technologies': [],
                'controls': [],
                'policies': []
            }
    
    def _extract_list(self, text: str, section_header: str) -> List[str]:
        """Extract list items from a section."""
        items = []
        
        # Find section
        if section_header not in text:
            return items
        
        section_start = text.find(section_header)
        section_text = text[section_start:]
        
        # Find next section or end
        next_section = section_text.find('\n\n')
        if next_section > 0:
            section_text = section_text[:next_section]
        
        # Extract items (lines starting with -)
        lines = section_text.split('\n')
        for line in lines:
            line = line.strip()
            if line.startswith('-'):
                item = line[1:].strip()
                if item:
                    items.append(item)
        
        return items
